import json
import re
from pptx import Presentation

def analyze (pptx_path, json_output_path=None):
    prs = Presentation(pptx_path)
    data = []

    for i, slide in enumerate(prs.slides):
        title = ''
        content = ''
        ntitle = 6  # Valor predeterminado alto (título más pequeño)
        has_links = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            content += text + '\n'

            # Detectar el título con base en el tamaño de fuente
            if shape.text_frame and shape.text_frame.paragraphs:
                font_sizes = [run.font.size.pt if run.font.size else 0
                              for para in shape.text_frame.paragraphs
                              for run in para.runs]
                max_size = max(font_sizes) if font_sizes else 0

                if max_size > 20 and max_size < ntitle * 10:
                    title = text
                    ntitle = int(70 // max_size) if max_size else 6

            # Verificar si hay links
            if re.search(r'https?://\S+', text):
                has_links = True

        slide_data = {
            "title": title,
            "content": content.strip(),
            "ntitle": ntitle,
            "links": has_links,
            "pages": i + 1
        }
        data.append(slide_data)

    # Guardar en archivo JSON if it was required
    if json_output_path:
        with open(json_output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=4)

    return data


analyze("press.pptx", "resultado.json")
# pip install python-pptx